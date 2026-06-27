"""check_layout.py - deterministic checks for the PowerPoint grader.

`extract_deck_digest()` is the only entry point. It returns the deck's id fields and
`computed_checks`: PASS/FAIL verdicts for the four rubric criteria that are a pure
data lookup the eye can't read reliably — font-family names, exact point sizes, and
hex colors. Everything visual (overlap, off-slide bleed, overflow, alignment, margins,
capitalization, source/page footers) is judged by the vision model from the rendered
slide IMAGES instead; those checks are deliberately not computed here.

The checks need no shape geometry, so extraction carries no box coordinates: it walks
each slide (recursing into groups) for the text-bearing shapes and their fonts, sizes,
colors, and placeholder role.
"""

import os
from collections import Counter
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError  # re-exported: grade_pptx.py catches it

__all__ = ["extract_deck_digest", "PackageNotFoundError"]

EMU_PER_INCH = 914400

# python-pptx placeholder-type enum name -> the short role tag logical_level reads.
_PH_ROLE = {
    "TITLE": "title", "CENTER_TITLE": "title", "BODY": "body", "SUBTITLE": "subtitle",
    "DATE": "dt", "FOOTER": "ftr", "SLIDE_NUMBER": "sldnum", "OBJECT": "obj",
}
@dataclass
class TextShape:
    # one text-bearing shape, reduced to what the deterministic checks read.
    name: str
    ph_type: str | None
    font_size_pt: float | None
    fonts: dict = field(default_factory=dict)    # family -> run count
    colors: list = field(default_factory=list)   # distinct run hexes seen


def _collect(shape, out: list):
    # append a TextShape for each text-bearing shape, recursing into groups
    # (geometry is irrelevant here). A run's color is its explicit RGB hex, or
    # skipped when it inherits the theme; the placeholder role maps the enum name
    # to a short tag for logical_level.
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect(child, out)
        return
    if not getattr(shape, "has_text_frame", False) or not shape.text_frame.text.strip():
        return

    fonts, sizes, colors = {}, {}, []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                fonts[run.font.name] = fonts.get(run.font.name, 0) + 1
            if run.font.size is not None:
                sizes[run.font.size.pt] = sizes.get(run.font.size.pt, 0) + 1
            c = run.font.color
            if c is not None and c.type is not None and c.rgb is not None and f"#{c.rgb}" not in colors:
                colors.append(f"#{c.rgb}")

    role = None
    if getattr(shape, "is_placeholder", False):
        pf = shape.placeholder_format
        role = _PH_ROLE.get(pf.type.name) if pf.type is not None else None

    out.append(TextShape(
        name=shape.name,
        ph_type=role,
        font_size_pt=max(sizes, key=sizes.get) if sizes else None,
        fonts=fonts,
        colors=colors,
    ))


def logical_level(shape: TextShape):
    # classify a text shape as title / body / footnote from its role and name.
    role = (shape.ph_type or "").lower()
    name = (shape.name or "").lower()
    if role == "title" or "title" in name:
        return "title"
    if role in ("ftr", "sldnum", "dt") or any(
            k in name for k in ("footer", "source", "footnote", "page", "ticker", "caption")):
        return "footnote"
    if role in ("body", "subtitle", "obj") or any(
            k in name for k in ("body", "subtitle", "bullet")):
        return "body"
    return None


# --------------------------------------------------------------------------- #
# The four deterministic checks. Each returns {verdict, detail}.
# --------------------------------------------------------------------------- #
def check_single_font_family(slides) -> dict:
    # at most 2 distinct font-family names across every run in the deck.
    families = sorted({fam for shapes in slides for s in shapes for fam in s.fonts})
    return {
        "verdict": "PASS" if len(families) <= 2 else "FAIL",
        "detail": f"{len(families)} distinct font families: {families}",
    }


def check_size_hierarchy(slides) -> dict:
    # per slide, title >= body + 4pt and body >= footnote + 2pt (when both exist).
    fails = []
    for index, shapes in enumerate(slides, start=1):
        per = {"title": [], "body": [], "footnote": []}
        for s in shapes:
            lvl = logical_level(s)
            if lvl and s.font_size_pt is not None:
                per[lvl].append(s.font_size_pt)
        title = max(per["title"]) if per["title"] else None
        body_min = min(per["body"]) if per["body"] else None
        body_max = max(per["body"]) if per["body"] else None
        foot_max = max(per["footnote"]) if per["footnote"] else None
        if title is not None and body_max is not None and title < body_max + 4:
            fails.append(f"slide {index}: title {title}pt not >= body {body_max}pt +4")
        if body_min is not None and foot_max is not None and body_min < foot_max + 2:
            fails.append(f"slide {index}: body {body_min}pt not >= footnote {foot_max}pt +2")
    return {
        "verdict": "FAIL" if fails else "PASS",
        "detail": "; ".join(fails) or "title>body>footnote size hierarchy holds on every slide",
    }


def check_consistent_level_sizes(slides) -> dict:
    # each logical level uses one point size deck-wide (at most 1 deviating run).
    level_sizes = {"title": [], "body": [], "footnote": []}
    for shapes in slides:
        for s in shapes:
            lvl = logical_level(s)
            if lvl and s.font_size_pt is not None:
                level_sizes[lvl].append(s.font_size_pt)
    details, fail = [], False
    for lvl, sizes in level_sizes.items():
        if not sizes:
            continue
        dominant, dom_n = Counter(sizes).most_common(1)[0]
        outliers = len(sizes) - dom_n
        details.append(f"{lvl}: dominant {dominant}pt, {outliers} outlier run(s)")
        if outliers >= 2:
            fail = True
    return {
        "verdict": "FAIL" if fail else "PASS",
        "detail": "; ".join(details) or "no leveled text found",
    }


def check_limited_palette(slides) -> dict:
    # at most 6 distinct non-neutral text hues across the deck. A pure gray
    # (R==G==B, so also black and white) is neutral, not a brand hue.
    hues = set()
    for shapes in slides:
        for s in shapes:
            for hx in s.colors:
                h = hx.lstrip("#").upper()[:6]
                if len(h) == 6 and not h[0:2] == h[2:4] == h[4:6]:
                    hues.add(h)
    return {
        "verdict": "PASS" if len(hues) <= 6 else "FAIL",
        "detail": f"{len(hues)} distinct non-neutral text hues: {sorted(hues)}",
    }


def extract_deck_digest(path: str) -> dict:
    # main entry: id fields + the four deterministic PASS/FAIL checks, keyed by
    # rubric criterion id so the grader can adopt each verdict as fact.
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            _collect(shape, shapes)
        slides.append(shapes)
    return {
        "file": os.path.basename(path),
        "slide_count": len(prs.slides),
        "slide_size_in": [
            round(prs.slide_width / EMU_PER_INCH, 3),
            round(prs.slide_height / EMU_PER_INCH, 3),
        ],
        "computed_checks": {
            "typography-single-font-family": check_single_font_family(slides),
            "typography-size-hierarchy": check_size_hierarchy(slides),
            "typography-consistent-level-sizes": check_consistent_level_sizes(slides),
            "color-and-branding-limited-palette": check_limited_palette(slides),
        },
    }
